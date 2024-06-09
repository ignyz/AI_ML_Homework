# Copyright (c) MaliciaLab, 2023.
# This code is licensed under the MIT license. 
# See the LICENSE file in the iocsearcher project root for license terms. 
#
import os
import logging
from io import BytesIO

import PyPDF2
from bs4 import BeautifulSoup

from pptx import Presentation
from docx import Document as DocxDocument
from iocsearcher.doc_base import Document
from iocsearcher.doc_epdf import ExtendedPdf
from iocsearcher.doc_ehtml import ExtendedHtml
from iocsearcher.doc_common import get_file_mime_type
from fastapi import FastAPI, File, UploadFile, HTTPException
from email import policy, message_from_bytes
from email.parser import BytesParser
import magic
import pandas as pd

# Set logging
log = logging.getLogger(__name__)
def get_mime_type(byte_stream):
    # Create a magic instance
    mime = magic.Magic(mime=True)
    # Get the MIME type
    mime_type = mime.from_buffer(byte_stream.read())
    # Reset the stream position to the beginning
    byte_stream.seek(0)
    return mime_type

def open_document(filepath, create_ioc_fun=None):
    """Return a Document object, None if unsupported document type"""
    # Check we have a file
    if not os.path.isfile(filepath):
        log.warning("Not a file: %s" % filepath)
        return None
    # Get MIME type
    mime_type = get_file_mime_type(filepath)
    log.debug("  File MIME type: %s" % mime_type)
    # Create right object according to MIME type
    tokens = mime_type.split('/')
    if tokens[1] == "pdf":
        doc = ExtendedPdf(filepath, mime_type=mime_type,
                          create_ioc_fun=create_ioc_fun)
    elif ((mime_type == "text/html") or
          (mime_type == "text/xml")):
          doc = ExtendedHtml(filepath, mime_type=mime_type,
                              create_ioc_fun=create_ioc_fun)
    elif ((tokens[0] == "text") or
          (mime_type == "application/csv") or
          (mime_type == "application/json")):
        doc = Document(filepath, mime_type=mime_type)
    else:
        log.warning("Unsupported MIME type %s for %s" % (mime_type, filepath))
        doc = None
    return doc
def read_presentation(byte_stream):
    presentation = Presentation(byte_stream)
    content = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        content.append("\n".join(slide_text))
    return "\n\n".join(content)

def read_word_doc(byte_stream):
    doc = DocxDocument(byte_stream)
    content = []
    for para in doc.paragraphs:
        content.append(para.text)
    return "\n".join(content)
def extract_text_from_html(html_file):
    soup = BeautifulSoup(html_file.decode('utf-8'), "lxml")
    # Extract text from all elements
    text = soup.get_text(separator="\n", strip=True)
    return text
def read_excel(byte_stream):
    # Assuming the Excel file has the first sheet as the target
    df = pd.read_excel(byte_stream)
    return df.to_string(index=False)

def open_document_bn(file, doc_content, create_ioc_fun=None):
    """Return a Document object, "" if unsupported document type"""
    # Get MIME type
    mime_type = file.content_type #get_file_mime_type(file)
    log.debug("  File MIME type: %s" % mime_type)

    # Create right object according to MIME type
    tokens = mime_type.split('/')
    if tokens[1] == "pdf":
        # Read pdf content
        try:
            # Create a PyPDF2 PdfFileReader object
            pdf_reader = PyPDF2.PdfReader(BytesIO(doc_content))
            doc = ''.join([pdf_reader.pages[page_num].extract_text() for page_num in range(len(pdf_reader.pages))])
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error processing PDF: {str(e)}")
    elif ((mime_type == "text/html") or
          (mime_type == "text/xml")):
          doc = extract_text_from_html(doc_content)
    elif ((tokens[0] == "text") or
        (mime_type == "application/csv") or
        (mime_type == "application/json")):
        doc = doc_content.decode('utf-8')
    elif mime_type == "message/rfc822":
        # Parse the email message from the BytesIO stream
        msg = message_from_bytes(doc_content)
        # Access email fields
        print("Subject:", msg['subject'])
        print("From:", msg['from'])
        print("To:", msg['to'])
        print("Date:", msg['date'])
        return str(msg)
    elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        return read_presentation(BytesIO(doc_content))
    elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        return read_word_doc(BytesIO(doc_content))
    elif mime_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'application/vnd.ms-excel']:
        return read_excel(BytesIO(doc_content))
    elif mime_type == 'application/x-msdownload':
        decoded_strings = ""  # Initialize an empty list to store decoded strings
        for chunk in doc_content.split(b'\0'):  # Split the content by null bytes
            try:
                decoded_string = chunk.decode('utf-8', 'ignore')  # Decode the chunk into a string, ignoring errors
                if decoded_string  != "":
                    decoded_strings += decoded_string   # Add the decoded string to the list
            except Exception as e:
                log.warning("Could not decode document content: %s" % str(e))
        return decoded_strings
    else:
        log.warning("Unsupported MIME type %s for %s" % (mime_type, file))
        doc = ""
    return doc

